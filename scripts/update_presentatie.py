from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation("20260603_selectietool_mvp_voorstel_intern.pptx")

WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(51, 51, 51)
GRAY = RGBColor(120, 120, 120)
BLUE = RGBColor(44, 62, 80)
ACCENT = RGBColor(41, 128, 185)


def add_text_slide(title_text, body_lines, note=None, insert_at=None):
    if insert_at is not None:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        # Move slide to correct position
        slides_el = prs.slides._sldIdLst
        slide_id_list = list(slides_el)
        last = slide_id_list[-1]
        slides_el.remove(last)
        slides_el.insert(insert_at, last)
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    title = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.9)
    )
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BLUE

    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.2))
    tf = body.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

        if line == "":
            p.space_before = Pt(6)
            continue

        bold = False
        if line.startswith("**") and line.endswith("**"):
            line = line[2:-2]
            bold = True

        p.text = line
        if bold:
            p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
        p.space_before = Pt(3)
        if line.startswith("- "):
            p.level = 0

    if note:
        nb = slide.shapes.add_textbox(
            Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.6)
        )
        tf = nb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = note
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY

    return slide


def add_table_slide(title_text, headers, rows, intro=None, note=None, insert_at=None):
    if insert_at is not None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slides_el = prs.slides._sldIdLst
        slide_id_list = list(slides_el)
        last = slide_id_list[-1]
        slides_el.remove(last)
        slides_el.insert(insert_at, last)
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    title = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.9)
    )
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BLUE

    top = Inches(1.4)
    if intro:
        ib = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.5)
        )
        tf = ib.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = intro
        p.font.size = Pt(15)
        p.font.color.rgb = DARK
        top = Inches(2.0)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_width = Inches(11.7)
    tbl_height = Inches(0.35) * n_rows
    shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.8), top, tbl_width, tbl_height
    )
    table = shape.table

    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = DARK

    if note:
        y = top + tbl_height + Inches(0.2)
        nb = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.7), Inches(0.5))
        tf = nb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = note
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY

    return slide


# ============================================================
# UPDATE SLIDE 3: Add old Psychologie data to comparison table
# ============================================================
slide3 = prs.slides[2]
for shape in slide3.shapes:
    if shape.has_table:
        table = shape.table
        # Add 2 new columns for old Psychologie data
        # Since python-pptx can't easily add columns, we'll update the intro text instead
        break

# Update subtitle to mention 5 bestanden
for shape in slide3.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            if "drie echte selectiebestanden" in para.text:
                for run in para.runs:
                    if "drie" in run.text:
                        run.text = run.text.replace("drie", "vijf")


# ============================================================
# ADD SLIDE: Old Psychologie data (insert after slide 14, before slide 15)
# ============================================================
add_text_slide(
    "MVP-keuzes: Psychologie 2021-2022 en 2022-2023 (9 kolommen)",
    [
        "Naast het uitgebreide Psychologie-bestand van 2026-2027 hebben we twee oudere bestanden. "
        "Deze zijn veel eenvoudiger: 9 kolommen elk, met al geaggregeerde scores.",
        "",
        "**Drie instrumenten per bestand**",
        "Selectietoets (ruwe score + percentage), Matchingsvragenlijst (ruwe score + percentage) en Cijferlijst (ruwe score + percentage). "
        "Plus een totaalscore en rangnummer.",
        "",
        "**Kolomnamen verschillen per jaar**",
        'De kolommen heten net anders: 2021-2022 heeft "TOETS", "MATCH", "CIJF" terwijl 2022-2023 "Toets", "Matching", "Cijferlijst" gebruikt. '
        "Per bestand is daarom een eigen config nodig.",
        "",
        "**Geen individuele vakscores**",
        "In tegenstelling tot 2026-2027 zijn er geen scores per schoolvak. De cijferlijst is al samengevat tot een enkel getal. "
        "Analyse per vak is dus niet mogelijk bij deze jaargangen.",
        "",
        "**Oudere bestandsformaat**",
        "Het 2021-2022 bestand is een .xls (oud Excel-formaat). De tool moet beide formaten aankunnen (.xls en .xlsx).",
    ],
    note="Per bestand gebruiken we 6 van de 9 kolommen. Studentnummer, Totaalscore en Rangnummer staan in de instellingen.",
    insert_at=14,
)


# ============================================================
# ADD SLIDE: Scope analysis - original spec vs MVP
# ============================================================
add_table_slide(
    "Oorspronkelijke opdracht vs. onze MVP",
    ["Onderdeel uit oorspronkelijke opdracht", "In MVP?", "Toelichting"],
    [
        [
            "Selectiedata uploaden en koppelen aan items",
            "Ja",
            "Kernfunctionaliteit van de tool",
        ],
        [
            "Drie variabelecategorieen: selectie, uitkomst, achtergrond",
            "Deels",
            "Selectievariabelen via config, uitkomst via 1CHO, achtergrond niet",
        ],
        [
            "Dashboard met groepsvergelijking per item",
            "Ja",
            "Drie groepen: niet geselecteerd, niet doorgestroomd, doorgestroomd",
        ],
        ["Koppeling met studiesuccesdata (1CHO)", "Ja", "Via studentnummer"],
        [
            "Beschrijvende statistiek per variabele",
            "Ja",
            "Gemiddelde, spreiding per groep in het dashboard",
        ],
        [
            "Correlatie- en regressieanalyse",
            "Nee",
            "Step-wise regressie, R-squared, p-waarden zijn post-MVP",
        ],
        [
            "Automatisch gegenereerd evaluatierapport (Word/PDF)",
            "Nee",
            "MVP toont een dashboard, geen downloadbaar rapport",
        ],
        [
            "Automatische tekstuele duiding van resultaten",
            "Nee",
            "Dashboard laat patronen zien, gebruiker interpreteert zelf",
        ],
        [
            "Interactieve variabele-definitie in de browser",
            "Nee",
            "We gebruiken een config-bestand (Excel). Sneller, betrouwbaarder voor MVP",
        ],
        [
            "Webapplicatie met authenticatie, rollen, 2FA",
            "Nee",
            "MVP draait lokaal of als eenvoudig dashboard",
        ],
        [
            "Dataveiligheid: pseudonimisering, verwijdertermijnen",
            "Nee",
            "Belangrijk, maar post-MVP",
        ],
        [
            "Achtergrondvariabelen en fairness-analyse",
            "Nee",
            "Uitsplitsing op geslacht, achtergrond etc. is post-MVP",
        ],
        [
            "Variabelen samenvoegen (bijv. synoniemen schoolvakken)",
            "Nee",
            "Tool neemt kolommen 1-op-1 over uit config",
        ],
        [
            "LOCS/HOCS/toepassing/communicatie categorisering",
            "Nee",
            "Toetskwaliteitsanalyse is post-MVP",
        ],
        [
            "Herkansingen vs. originele toetsen onderscheiden",
            "Nee",
            "Niet relevant bij selectiedata (alleen bij studieresultaten)",
        ],
        ["CSV-ondersteuning", "Nee", "MVP werkt alleen met Excel (.xlsx en .xls)"],
        ["SurfConext login", "Nee", "Geen multi-user authenticatie in MVP"],
        [
            "Config-bestand automatisch genereren via tool",
            "Nee",
            "Handmatig invullen via template. Automatisering is volgende stap",
        ],
    ],
    intro="De oorspronkelijke opdrachtomschrijving (NRO, 2023) en de technische ontwerpen (Dialogic, Yard) beschrijven een complete webapplicatie. Onze MVP pakt het analysedeel.",
    note='Bron: "Evaluatietool Selectie Hoger Onderwijs" (NRO 2023), "Ontwerpfase ETHO" (Dialogic 2023), "Technisch ontwerp" (Yard 2023)',
    insert_at=16,
)

# ============================================================
# ADD SLIDE: More scope items - what's uncertain
# ============================================================
add_text_slide(
    "Scope: twijfelgevallen en open vragen",
    [
        "Bij een aantal onderdelen uit de oorspronkelijke opdracht is het de vraag of en wanneer we ze meenemen.",
        "",
        "**Cohortanalyse (meerdere jaargangen vergelijken)**",
        "De oorspronkelijke opdracht vraagt om analyse per cohort en over cohorten heen. We hebben nu data van 2021 t/m 2026. "
        "Het dashboard kan per jaargang draaien, maar vergelijking tussen jaargangen is complexer. Opnemen als de basisversie werkt?",
        "",
        "**Gemiddelden berekenen over groepen variabelen**",
        "De opdracht noemt: gemiddelde van alle LOCS-resultaten, alle toetsen per type, alle vwo-cijfers. "
        'Dit vereist dat de config een extra classificatie krijgt (bijv. "LOCS", "HOCS"). Niet in MVP, maar de config-structuur kan het later aan.',
        "",
        "**Meerdere bestanden per selectieprocedure**",
        "De opdracht beschrijft dat data van een kandidaat in meerdere bestanden kan zitten. "
        "In de MVP laden we per keer een selectiebestand. Meerdere bestanden combineren is mogelijk maar nog niet gebouwd.",
        "",
        "**Toetskwaliteitsanalyse (item-niveau)**",
        "De opdracht noemt als toekomstige uitbreiding: analyse van individuele toetsvragen. "
        "Daar hebben we nu geen data voor, dus dit is sowieso post-MVP.",
    ],
    insert_at=17,
)


# ============================================================
# UPDATE SLIDE 17 (now 20): samenvattingstabel - add new data
# ============================================================
# Find the summary table slide (originally 17, now shifted)
# We need to find it by content since we inserted slides
for idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            rows = []
            for row in table.rows:
                rows.append([cell.text.strip() for cell in row.cells])
            if (
                len(rows) > 1
                and rows[0][0] == "Bestand"
                and "Kolommen totaal" in rows[0][1]
            ):
                # Add rows for new Psychologie data
                # We can't easily add rows to existing table, so update intro text
                for shape2 in slide.shapes:
                    if shape2.has_text_frame:
                        for para in shape2.text_frame.paragraphs:
                            if "klein deel" in para.text:
                                for run in para.runs:
                                    run.text = run.text.replace(
                                        "Per bestand gebruiken we een klein deel van de kolommen. De rest is bewust buiten scope.",
                                        "Per bestand gebruiken we een klein deel van de kolommen. De rest is bewust buiten scope. De oudere Psychologie-bestanden (2021-2022 en 2022-2023) hebben elk 9 kolommen waarvan we er 6 gebruiken.",
                                    )
                break


# ============================================================
# UPDATE SLIDE 18 (now shifted): volgende stappen - add config tool note
# ============================================================
for idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            full_text = " ".join(p.text for p in shape.text_frame.paragraphs)
            if "Klein beginnen" in full_text and "dashboard" in full_text:
                tf = shape.text_frame
                p = tf.add_paragraph()
                p.text = ""
                p.space_before = Pt(8)
                p = tf.add_paragraph()
                p.text = "Toekomstig: config-tool"
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = DARK
                p = tf.add_paragraph()
                p.text = "Op termijn kan een tool het aanmaken van configuratiebestanden vereenvoudigen: kolommen uit het selectiebestand tonen en de analist laten koppelen via een interface. Dat is voor nu out-of-scope, maar de config-structuur is er al op voorbereid."
                p.font.size = Pt(16)
                p.font.color.rgb = DARK
                break


prs.save("20260603_selectietool_mvp_voorstel_intern.pptx")
print("Presentatie opgeslagen.")
