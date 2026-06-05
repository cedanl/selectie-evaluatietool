import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation("20260603_selectietool_mvp_voorstel_intern.pptx")

WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(51, 51, 51)
GRAY = RGBColor(120, 120, 120)
BLUE = RGBColor(44, 62, 80)

fixes = []


def find_and_replace_text(slide, old, new):
    """Replace text in all shapes on a slide."""
    count = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if old in run.text:
                        run.text = run.text.replace(old, new)
                        count += 1
    return count


# ============================================================
# SLIDE 7: Remove rangnummer reference from config description
# ============================================================
slide7 = prs.slides[6]
for shape in slide7.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if (
                    "Op welke rij staan de kolomnamen? Welke kolom is de totaalscore? En het rangnummer?"
                    in run.text
                ):
                    run.text = run.text.replace(
                        "Op welke rij staan de kolomnamen? Welke kolom is de totaalscore? En het rangnummer?",
                        "Welke kolom is de totaalscore?",
                    )
                    fixes.append(
                        "Slide 7: rangnummer referentie verwijderd uit settings beschrijving"
                    )

# ============================================================
# SLIDE 8: Fix blad_naam value (remove "default eerste blad" parenthetical)
# ============================================================
slide8 = prs.slides[7]
for shape in slide8.shapes:
    if shape.has_table:
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if cells[0] == "blad_naam" and "(default eerste blad)" in cells[1]:
                # Keep the note but in a cleaner way - actually just use the real value
                # The PPT already has this intentionally, so we leave it
                pass

# ============================================================
# SLIDE 11: "drie" -> "vijf"
# ============================================================
n = find_and_replace_text(
    prs.slides[10], "de drie aangeleverde bestanden", "de vijf aangeleverde bestanden"
)
if n:
    fixes.append("Slide 11: drie -> vijf bestanden")

# ============================================================
# SLIDE 16: Be explicit about local
# ============================================================
slide16 = prs.slides[15]
for shape in slide16.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "De tool herkent kolommen niet zelf" in run.text:
                    run.text = run.text.replace(
                        "De tool herkent kolommen niet zelf. De analist maakt per bestand een configuratiebestand. Dat kost eenmalig werk, maar voorkomt fouten.",
                        "De tool herkent kolommen niet zelf. De analist maakt per bestand een configuratiebestand. Dat kost eenmalig werk, maar voorkomt fouten. Los van de oorspronkelijke opdracht (die een interactieve browser-interface beschrijft) kiezen wij bewust voor een Excel-config omdat het sneller is en de analist de data het beste kent.",
                    )
                    fixes.append(
                        "Slide 16: config-keuze losgekoppeld van Dialogic/Yard"
                    )

# ============================================================
# SLIDE 17: Replace the huge table with two smaller slides
# Delete current slide 17 and insert two new ones
# ============================================================
# We can't easily delete a slide in python-pptx, so we'll clear it and rebuild
slide17 = prs.slides[16]

# Clear all shapes
for shape in list(slide17.shapes):
    sp = shape._element
    sp.getparent().remove(sp)

slide17.background.fill.solid()
slide17.background.fill.fore_color.rgb = WHITE

# Rebuild slide 17 as "What we DO"
title = slide17.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.9))
tf = title.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Oorspronkelijke opdracht vs. onze MVP: wat we WEL doen"
p.font.size = Pt(26)
p.font.bold = True
p.font.color.rgb = BLUE

intro = slide17.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.5))
tf = intro.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "De oorspronkelijke opdracht (NRO 2023) en technische ontwerpen (Dialogic, Yard) beschrijven een complete webapplicatie. Onze MVP pakt het analysedeel en draait lokaal."
p.font.size = Pt(14)
p.font.color.rgb = DARK

n_rows = 7
n_cols = 3
shape = slide17.shapes.add_table(
    n_rows, n_cols, Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.38) * n_rows
)
table = shape.table
table.columns[0].width = Inches(5.0)
table.columns[1].width = Inches(1.2)
table.columns[2].width = Inches(5.5)

headers = ["Onderdeel", "MVP?", "Toelichting"]
yes_rows = [
    ["Selectiedata uploaden en koppelen aan items", "Ja", "Kernfunctionaliteit"],
    [
        "Dashboard met groepsvergelijking per item",
        "Ja",
        "Drie groepen: niet geselecteerd / niet doorgestroomd / doorgestroomd",
    ],
    ["Koppeling met studiesuccesdata (1CHO)", "Ja", "Via studentnummer"],
    [
        "Beschrijvende statistiek per variabele",
        "Ja",
        "Gemiddelde, spreiding per groep in het dashboard",
    ],
    [
        "Selectie- en uitkomstvariabelen",
        "Ja",
        "Selectievariabelen via config, uitkomst via 1CHO",
    ],
    [
        "Lokaal draaiend dashboard (Python)",
        "Ja",
        "Geen webapplicatie, geen hosting, geen login nodig",
    ],
]

for ci, h in enumerate(headers):
    cell = table.cell(0, ci)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
    cell.fill.solid()
    cell.fill.fore_color.rgb = BLUE

for ri, row in enumerate(yes_rows):
    for ci, val in enumerate(row):
        cell = table.cell(ri + 1, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.color.rgb = DARK

fixes.append('Slide 17: grote tabel gesplitst in "wat we wel doen"')


# ============================================================
# SLIDE 18: Rebuild as "what we DON'T do"
# ============================================================
slide18 = prs.slides[17]
for shape in list(slide18.shapes):
    sp = shape._element
    sp.getparent().remove(sp)

slide18.background.fill.solid()
slide18.background.fill.fore_color.rgb = WHITE

title = slide18.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.9))
tf = title.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Oorspronkelijke opdracht vs. onze MVP: wat we NIET doen"
p.font.size = Pt(26)
p.font.bold = True
p.font.color.rgb = BLUE

intro = slide18.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.5))
tf = intro.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Deze onderdelen uit de oorspronkelijke opdracht vallen buiten de MVP. Ze kunnen later worden toegevoegd."
p.font.size = Pt(14)
p.font.color.rgb = DARK

n_rows = 13
n_cols = 2
shape = slide18.shapes.add_table(
    n_rows, n_cols, Inches(0.8), Inches(1.9), Inches(11.7), Inches(0.35) * n_rows
)
table = shape.table
table.columns[0].width = Inches(5.0)
table.columns[1].width = Inches(6.7)

no_headers = ["Onderdeel", "Waarom niet in MVP"]
no_rows = [
    [
        "Correlatie- en regressieanalyse",
        "Step-wise regressie, R-squared, p-waarden: post-MVP",
    ],
    [
        "Evaluatierapport genereren (Word/PDF)",
        "MVP toont een dashboard, geen downloadbaar rapport",
    ],
    [
        "Automatische tekstuele duiding",
        "Dashboard laat patronen zien, gebruiker interpreteert zelf",
    ],
    [
        "Variabele-definitie via browser-interface",
        "Wij kiezen bewust voor een Excel-config: sneller en betrouwbaarder",
    ],
    [
        "Webapplicatie met login, rollen, 2FA",
        "MVP draait lokaal op de machine van de analist",
    ],
    [
        "Pseudonimisering en verwijdertermijnen",
        "Data blijft lokaal bij de analist, geen hosting",
    ],
    [
        "Achtergrondvariabelen en fairness-analyse",
        "Uitsplitsing op geslacht, achtergrond etc. is post-MVP",
    ],
    [
        "Variabelen samenvoegen (synoniemen)",
        "Tool neemt kolommen 1-op-1 over uit config",
    ],
    [
        "Categorisering van toetsvaardigheden",
        "Indeling in cognitieve niveaus (kennis, toepassing, analyse) is post-MVP",
    ],
    ["CSV-ondersteuning", "MVP werkt alleen met Excel (.xlsx en .xls)"],
    ["SurfConext login", "Geen multi-user authenticatie in MVP"],
    [
        "Config-tool (automatisch genereren)",
        "Handmatig invullen via template. Automatisering is volgende stap",
    ],
]

for ci, h in enumerate(no_headers):
    cell = table.cell(0, ci)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
    cell.fill.solid()
    cell.fill.fore_color.rgb = BLUE

for ri, row in enumerate(no_rows):
    for ci, val in enumerate(row):
        cell = table.cell(ri + 1, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.color.rgb = DARK

note = slide18.shapes.add_textbox(Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.5))
tf = note.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = 'Bron: "Evaluatietool Selectie Hoger Onderwijs" (NRO 2023), "Ontwerpfase ETHO" (Dialogic 2023), "Technisch ontwerp" (Yard 2023)'
p.font.size = Pt(10)
p.font.color.rgb = GRAY

fixes.append('Slide 18: grote tabel gesplitst in "wat we niet doen" (geen LOCS jargon)')


# ============================================================
# SLIDE 19 (was 18 - twijfelgevallen): Fix cohort text and LOCS
# ============================================================
# This slide was the old "twijfelgevallen" slide, now at index 18
# Actually after our rebuild of 17 and 18, the old slide 18 content
# moved. Let me find it by content.
for idx, slide in enumerate(prs.slides):
    found = False
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if "twijfelgevallen" in para.text.lower():
                    found = True
                    break
        if found:
            break
    if not found:
        continue

    # Found the twijfelgevallen slide
    # Clear and rebuild
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    title = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.9)
    )
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Open vragen en twijfelgevallen"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BLUE

    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.2))
    tf = body.text_frame
    tf.word_wrap = True

    lines = [
        ("Cohortanalyse: vergelijken over jaargangen", True),
        (
            "We hebben data van 2021 t/m 2026, maar de selectieprocedure verandert regelmatig. "
            "Instrumenten worden aangepast, weggehaald of vervangen. In de huidige dynamische omgeving "
            "van selectie in het hoger onderwijs is het niet vanzelfsprekend dat je jaargangen naast "
            "elkaar kunt leggen. Het dashboard kan per jaargang draaien, maar directe vergelijking "
            "tussen jaargangen vraagt om extra aandacht voor wat er is veranderd.",
            False,
        ),
        ("", False),
        ("Meerdere bestanden per selectieprocedure", True),
        (
            "De oorspronkelijke opdracht beschrijft dat data van een kandidaat in meerdere bestanden kan zitten. "
            "In de MVP laden we per keer een selectiebestand. Meerdere bestanden combineren is mogelijk maar nog niet gebouwd.",
            False,
        ),
        ("", False),
        ("Gemiddelden over groepen variabelen", True),
        (
            "De opdracht noemt gemiddelden over categorieën toetsen (kennistoetsen, toepassingstoetsen, etc.). "
            "Dit vereist een extra classificatie in de config. Niet in MVP, maar de structuur kan het later aan.",
            False,
        ),
        ("", False),
        ("Toetskwaliteitsanalyse", True),
        (
            "Analyse van individuele toetsvragen, als toekomstige uitbreiding genoemd in de oorspronkelijke opdracht. "
            "We hebben daar nu geen data voor, dus sowieso post-MVP.",
            False,
        ),
    ]

    for i, (text, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if text == "":
            p.space_before = Pt(6)
            continue
        p.text = text
        p.font.bold = bold
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
        p.space_before = Pt(3)

    fixes.append("Slide (twijfelgevallen): cohort-tekst herschreven, LOCS verwijderd")
    break


# ============================================================
# SLIDE 20: Update samenvattingstabel - add 2 new Psychologie rows
# ============================================================
for idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            if (
                len(rows) > 1
                and rows[0][0] == "Bestand"
                and "Kolommen totaal" in rows[0][1]
            ):
                # Clear and rebuild this slide with 5-row table
                for s in list(slide.shapes):
                    sp = s._element
                    sp.getparent().remove(sp)

                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = WHITE

                title = slide.shapes.add_textbox(
                    Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.9)
                )
                tf = title.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = "Wat we meenemen en wat we overslaan"
                p.font.size = Pt(28)
                p.font.bold = True
                p.font.color.rgb = BLUE

                intro = slide.shapes.add_textbox(
                    Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.5)
                )
                tf = intro.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = "Per bestand gebruiken we een klein deel van de kolommen. De rest is bewust buiten scope."
                p.font.size = Pt(15)
                p.font.color.rgb = DARK

                headers = [
                    "Bestand",
                    "Kolommen totaal",
                    "In config",
                    "Overgeslagen",
                    "Voornaamste reden",
                ]
                data_rows = [
                    [
                        "FAR Leiden 2025",
                        "60",
                        "11",
                        "49",
                        "Dubbele beoordelaars, tekstvelden, persoonsgegevens",
                    ],
                    [
                        "FAR Leiden 2026",
                        "97",
                        "12",
                        "85",
                        "Drie scoreversies, proctoring, procesdata",
                    ],
                    [
                        "Psychologie 2026-2027",
                        "46",
                        "19",
                        "27",
                        "Vaknamen, ruwe cijfers, herhaalde structuur",
                    ],
                    [
                        "Psychologie 2022-2023",
                        "9",
                        "6",
                        "3",
                        "Eenvoudig bestand, alleen ID + totaalscore + rang over",
                    ],
                    [
                        "Psychologie 2021-2022",
                        "9",
                        "6",
                        "3",
                        "Zelfde structuur als 2022-2023, andere kolomnamen (.xls)",
                    ],
                ]

                n_rows = len(data_rows) + 1
                shape = slide.shapes.add_table(
                    n_rows,
                    5,
                    Inches(0.8),
                    Inches(2.0),
                    Inches(11.7),
                    Inches(0.38) * n_rows,
                )
                tbl = shape.table

                for ci, h in enumerate(headers):
                    cell = tbl.cell(0, ci)
                    cell.text = h
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(12)
                        p.font.bold = True
                        p.font.color.rgb = WHITE
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = BLUE

                for ri, row in enumerate(data_rows):
                    for ci, val in enumerate(row):
                        cell = tbl.cell(ri + 1, ci)
                        cell.text = val
                        for p in cell.text_frame.paragraphs:
                            p.font.size = Pt(11)
                            p.font.color.rgb = DARK

                fixes.append("Slide 20: samenvattingstabel uitgebreid naar 5 bestanden")
                break
    if fixes and "Slide 20" in fixes[-1]:
        break


# ============================================================
# SLIDE 21: "drie" -> "vijf" in volgende stappen
# ============================================================
for idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            full = " ".join(p.text for p in shape.text_frame.paragraphs)
            if "Volgende stappen" in full and "Inleesroutine" in full:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "drie selectiebestanden" in run.text:
                            run.text = run.text.replace(
                                "drie selectiebestanden", "vijf selectiebestanden"
                            )
                            fixes.append("Slide 21: drie -> vijf selectiebestanden")
                        if "alle drie de bestanden" in run.text:
                            run.text = run.text.replace(
                                "alle drie de bestanden", "alle vijf de bestanden"
                            )
                            fixes.append("Slide 21: alle drie -> alle vijf")
                        # Make explicit it runs locally
                        if "Taal: Python" in run.text:
                            run.text = run.text.replace(
                                "Taal: Python.",
                                "Taal: Python. De tool draait lokaal op de machine van de analist.",
                            )
                            fixes.append("Slide 21: lokaal expliciet gemaakt")
                break


try:
    prs.save("20260603_selectietool_mvp_voorstel_intern.pptx")
except PermissionError:
    prs.save("20260604_selectietool_mvp_voorstel_intern.pptx")
    print("Opgeslagen als 20260604 (origineel is open in PowerPoint).")
print("Presentatie gefixt. %d fixes:" % len(fixes))
for f in fixes:
    print("  - " + f)
