import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation("20260604_selectietool_mvp_voorstel_intern.pptx")
DARK = RGBColor(51, 51, 51)
BLUE = RGBColor(44, 62, 80)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(120, 120, 120)


def clear_slide(slide):
    for shape in list(slide.shapes):
        slide.shapes._spTree.remove(shape._element)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE


def add_title(slide, text, size=28):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = BLUE


def add_body(slide, lines, top=1.4, fontsize=16):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.7), Inches(5.4))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if text == "":
            p.space_before = Pt(6)
            continue
        p.text = text
        p.font.bold = bold
        p.font.size = Pt(fontsize)
        p.font.color.rgb = DARK
        p.space_before = Pt(3)


# ============================================================
# 1. Slide 8: rewrite with LLM mention
# ============================================================
slide8 = prs.slides[7]
clear_slide(slide8)
add_title(slide8, "Het configuratiebestand: nu handmatig, later via de tool")
add_body(
    slide8,
    [
        (
            "In de MVP vult de analist het configuratiebestand handmatig in via een Excel-template. "
            "Er is een uitleg-tabblad en elke cel heeft een toelichting. Dat werkt, maar kost tijd.",
            False,
        ),
        ("", False),
        (
            "Als volgende stap kan de tool zelf helpen bij het aanmaken van de config:",
            True,
        ),
        ("", False),
        ("1. De analist uploadt het selectiebestand in de tool.", False),
        ("2. De tool toont alle kolommen uit het bestand.", False),
        ("3. Per kolom kiest de analist: meenemen of overslaan?", False),
        (
            "4. Voor kolommen die worden meegenomen, vult de analist instrument, item en criterium in.",
            False,
        ),
        (
            "5. De tool genereert de config en gebruikt die direct voor de analyse.",
            False,
        ),
        ("", False),
        (
            "Een verdere stap: een taalmodel (LLM) kan op basis van de kolomnamen en de eerste rijen data "
            "een voorstel doen voor de koppeling. De analist hoeft dan alleen te bevestigen of bij te sturen.",
            False,
        ),
        ("", False),
        (
            "Beide opties zijn out-of-scope voor de MVP. Eerst de basisflow werkend krijgen "
            "(inlezen, omzetten, dashboard), dan het config-stuk vereenvoudigen.",
            False,
        ),
    ],
)
print("Slide 8 bijgewerkt.")


# ============================================================
# 2. Slide 19 (NIET doen): fix two rows
# ============================================================
slide19 = prs.slides[18]
for shape in slide19.shapes:
    if not shape.has_table:
        continue
    table = shape.table
    for row in table.rows:
        c0 = row.cells[0].text.strip()
        if "Variabelen samenvoegen" in c0:
            row.cells[2].text = (
                "Schoolvakken heten in verschillende bestanden net anders "
                '(bijv. "Wiskunde B", "WISB", "Wiskunde-B"). De NRO-opdracht '
                "beschrijft een functie om synoniemen samen te voegen. "
                "In onze MVP neemt de analist kolommen 1-op-1 over via de config. "
                "Samenvoegen kan de analist zelf doen in de brondata."
            )
            for p in row.cells[2].text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = DARK
            print("Slide 19: variabelen samenvoegen uitgebreid.")

        if "variabele-definitie" in c0.lower() or "Volledig webportaal" in c0:
            row.cells[0].text = "Volledig webportaal voor variabele-definitie"
            for p in row.cells[0].text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = DARK
            row.cells[2].text = (
                "Yard ontwierp een complete webapplicatie met login, rollen, "
                "instituutsbeheer en database. Wij bouwen een lokale tool. "
                "De config-stap kan later via een interface in de tool zelf (zie slide 8)."
            )
            for p in row.cells[2].text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = DARK
            print("Slide 19: variabele-definitie herformuleerd.")


# ============================================================
# 3. New slide: argument against filter in 1cijferho tool
# ============================================================
new_slide = prs.slides.add_slide(prs.slide_layouts[6])
slides_el = prs.slides._sldIdLst
slide_ids = list(slides_el)
last = slide_ids[-1]
slides_el.remove(last)
slides_el.insert(20, last)

new_slide.background.fill.solid()
new_slide.background.fill.fore_color.rgb = WHITE
add_title(new_slide, "Filteren van 1CHO-data: waar hoort dat thuis?")
add_body(
    new_slide,
    [
        (
            "Er is voorgesteld om een filterfunctie toe te voegen aan de 1cijferho tool, "
            "zodat de analist daar al de juiste subset selecteert voordat de data naar "
            "de evaluatietool gaat.",
            False,
        ),
        ("", False),
        ("Dat is onnodig extra werk. Er zijn drie eenvoudigere alternatieven:", True),
        ("", False),
        ("1. De evaluatietool filtert zelf", True),
        (
            "De evaluatietool kan bij het inlezen van de 1CHO-data filteren op opleiding, "
            "cohort en relevante kolommen. Dat is een paar regels code in de tool "
            "die we toch al bouwen.",
            False,
        ),
        ("", False),
        ("2. De analist filtert in Excel of met code", True),
        (
            "De gebruikers zijn data professionals. Ze kunnen een Excel-bestand filteren "
            "voordat ze het uploaden. Dit kost minder dan een minuut en vereist geen "
            "aanpassing aan een andere tool.",
            False,
        ),
        ("", False),
        ("3. De 1cijferho tool heeft al presets", True),
        (
            "De 1cijferho tool heeft een preset-functie waarmee je alleen de relevante "
            "kolommen krijgt. De data is op kolomniveau dus al kleiner. Extra filteren "
            "op rijniveau past beter in de tool die de data ook daadwerkelijk gebruikt.",
            False,
        ),
        ("", False),
        (
            "De doelgroep is op dit moment twee opleidingen, met data professionals "
            "die al aangehaakt zijn. Een filterfunctie bouwen in de 1cijferho tool is "
            "bouwen voor een probleem dat er niet is. Als filteren later een bottleneck "
            "wordt, kunnen we het alsnog overwegen.",
            False,
        ),
    ],
    fontsize=15,
)
print("Nieuwe slide: 1cijferho filter argument.")


prs.save("20260604_selectietool_mvp_voorstel_intern.pptx")
print("Klaar.")
