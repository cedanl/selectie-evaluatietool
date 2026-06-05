import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation("20260604_selectietool_mvp_voorstel_intern.pptx")

WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(51, 51, 51)
GRAY = RGBColor(120, 120, 120)
BLUE = RGBColor(44, 62, 80)

fixes = []


def clear_slide(slide):
    for shape in list(slide.shapes):
        slide.shapes._spTree.remove(shape._element)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE


def add_title(slide, text, size=28):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = BLUE


def add_intro(slide, text, top=1.3):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.7), Inches(0.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = DARK


def add_note(slide, text, top=6.6):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.7), Inches(0.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY


def add_table(slide, headers, rows, top=2.0, row_height=0.36):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    shape = slide.shapes.add_table(
        n_rows,
        n_cols,
        Inches(0.8),
        Inches(top),
        Inches(11.7),
        Inches(row_height) * n_rows,
    )
    table = shape.table
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = DARK
    return shape


def add_body(slide, lines, top=1.5):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.7), Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if text == "":
            p.space_before = Pt(6)
            continue
        p.text = text
        p.font.bold = bold
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
        p.space_before = Pt(3)


# ============================================================
# INSERT NEW SLIDE after slide 7: Config-tool toekomst
# ============================================================
new_slide = prs.slides.add_slide(prs.slide_layouts[6])
slides_el = prs.slides._sldIdLst
slide_ids = list(slides_el)
last = slide_ids[-1]
slides_el.remove(last)
slides_el.insert(7, last)  # insert after slide 7 (0-indexed)

new_slide.background.fill.solid()
new_slide.background.fill.fore_color.rgb = WHITE
add_title(new_slide, "Waarom een config-bestand en geen tool?")
add_body(
    new_slide,
    [
        (
            "In de oorspronkelijke plannen (Yard, 2023) was het idee dat de analist kolommen koppelt via een "
            "interactieve webinterface: het systeem toont de kolommen uit het bestand, en de analist klikt "
            "per kolom aan wat het is.",
            False,
        ),
        ("", False),
        ("Wij kiezen in de MVP bewust voor een Excel-config. Drie redenen:", False),
        ("", False),
        ("1. Sneller te bouwen", True),
        (
            "Een Excel-template met uitleg is vandaag klaar. Een interactieve koppel-interface kost weken.",
            False,
        ),
        ("", False),
        ("2. De analist kent de data het beste", True),
        (
            "De analist weet welke kolommen relevant zijn en welke niet. Een Excel-bestand dwingt tot nadenken over die keuze.",
            False,
        ),
        ("", False),
        ("3. Herbruikbaar en controleerbaar", True),
        (
            "Een config-bestand kun je delen, versiebeheren en door een collega laten checken. Dat is lastiger met een interface die je doorklinkt.",
            False,
        ),
        ("", False),
        (
            "Op termijn kan een tool het aanmaken van configs vereenvoudigen: kolommen uit het bestand tonen, "
            "laten koppelen via drag-and-drop, en de config automatisch genereren. Maar dat is pas zinvol "
            "als de basisstructuur bewezen is.",
            False,
        ),
    ],
    top=1.4,
)
fixes.append("Nieuwe slide 8: config-tool vs Excel uitleg")


# ============================================================
# SLIDE 10 (dashboard, now at index 10 because we inserted one):
# Add boxplot, correlatie, regressie specifics
# ============================================================
# After inserting a slide, the old slide 10 is now at index 10
slide_dashboard = prs.slides[10]
clear_slide(slide_dashboard)
add_title(slide_dashboard, "Wat laat het dashboard zien?")
add_body(
    slide_dashboard,
    [
        (
            "Het dashboard toont per selectie-item een vergelijking tussen drie groepen kandidaten.",
            False,
        ),
        ("", False),
        ("Groep 1: Niet geselecteerd / niet begonnen", True),
        (
            "Kandidaten die onder de streep vielen of niet zijn gaan studeren (niet in het 1CHO bestand).",
            False,
        ),
        ("", False),
        ("Groep 2: Geselecteerd, niet doorgestroomd", True),
        (
            "Kandidaten die zijn toegelaten, maar niet naar jaar 2 zijn gegaan of geen diploma hebben gehaald.",
            False,
        ),
        ("", False),
        ("Groep 3: Geselecteerd, doorgestroomd", True),
        (
            "Kandidaten die zijn toegelaten en het goed hebben gedaan (doorstroom naar jaar 2 of diploma).",
            False,
        ),
        ("", False),
        (
            "In de MVP tonen we per item een boxplot van de scores per groep. Zo zie je in een oogopslag of de groepen van elkaar verschillen.",
            False,
        ),
        ("", False),
        (
            "Later kan het dashboard worden uitgebreid met correlatie-analyse (welke items hangen samen met studiesucces?) en regressie-analyse (welke items voorspellen studiesucces het best, en hoeveel voegt elk item toe?).",
            False,
        ),
    ],
    top=1.4,
)
fixes.append("Slide 11 (dashboard): boxplot, correlatie en regressie specifiek benoemd")


# ============================================================
# SLIDE 17 (now index 17): "WEL doen" - rename + add source + add CSV
# ============================================================
slide17 = prs.slides[17]
clear_slide(slide17)
add_title(slide17, "Oorspronkelijk idee vs. onze MVP: wat we WEL doen", size=26)
add_intro(
    slide17,
    "Het oorspronkelijke idee (NRO, 2023) en technische ontwerpen (Dialogic, Yard) beschrijven een complete webapplicatie. Onze MVP pakt het analysedeel en draait lokaal.",
)

add_table(
    slide17,
    ["Onderdeel", "Bron", "Toelichting"],
    [
        [
            "Selectiedata uploaden en koppelen aan items",
            "NRO, Yard",
            "Kernfunctionaliteit van de tool",
        ],
        [
            "Dashboard met boxplots per groep",
            "MVP",
            "Drie groepen: niet geselecteerd / niet doorgestroomd / doorgestroomd",
        ],
        ["Koppeling met studiesuccesdata (1CHO)", "NRO", "Via studentnummer"],
        [
            "Beschrijvende statistiek per variabele",
            "NRO, Dialogic",
            "Gemiddelde, spreiding per groep in het dashboard",
        ],
        [
            "Selectie- en uitkomstvariabelen",
            "NRO",
            "Selectievariabelen via config, uitkomst via 1CHO",
        ],
        [
            "Excel- en CSV-ondersteuning voor data",
            "NRO, Yard",
            "Databestanden in .xlsx, .xls of .csv. Config blijft Excel",
        ],
        [
            "Lokaal draaiend dashboard (Python)",
            "MVP",
            "Geen webapplicatie, geen hosting, geen login nodig",
        ],
    ],
)
fixes.append('Slide 18: "idee" i.p.v. "opdracht", bronnen toegevoegd, CSV als ja')


# ============================================================
# SLIDE 18 (now index 18): "NIET doen" - rename + add source + expand
# ============================================================
slide18 = prs.slides[18]
clear_slide(slide18)
add_title(slide18, "Oorspronkelijk idee vs. onze MVP: wat we NIET doen", size=26)
add_intro(
    slide18,
    "Deze onderdelen uit het oorspronkelijke idee vallen buiten de MVP. Ze kunnen later worden toegevoegd.",
)

add_table(
    slide18,
    ["Onderdeel", "Bron", "Waarom niet in MVP"],
    [
        [
            "Correlatie- en regressie-analyse",
            "NRO, Dialogic",
            "Step-wise regressie, R-squared, p-waarden: post-MVP. Boxplots eerst",
        ],
        [
            "Evaluatierapport genereren (Word/PDF)",
            "NRO, Dialogic",
            "MVP toont een dashboard, geen downloadbaar rapport",
        ],
        [
            "Automatische tekstuele duiding",
            "NRO, Dialogic",
            "Dashboard laat patronen zien, gebruiker interpreteert zelf",
        ],
        [
            "Interactieve variabele-definitie in de browser",
            "Yard",
            "Het idee was dat de analist kolommen koppelt via een webinterface. "
            "Wij kiezen voor een Excel-config: sneller, controleerbaar, deelbaar",
        ],
        [
            "Webapplicatie met login, rollen, 2FA",
            "Yard",
            "MVP draait lokaal op de machine van de analist",
        ],
        [
            "Pseudonimisering en verwijdertermijnen",
            "NRO, Yard",
            "Data blijft lokaal bij de analist, geen hosting",
        ],
        [
            "Achtergrondvariabelen en fairness-analyse",
            "NRO",
            "Uitsplitsing op geslacht, achtergrond etc. is post-MVP",
        ],
        [
            "Variabelen samenvoegen (synoniemen)",
            "NRO",
            "Tool neemt kolommen 1-op-1 over uit config",
        ],
        [
            "Categorisering van toetsvaardigheden",
            "NRO",
            "Indeling in cognitieve niveaus (kennis, toepassing, analyse) is post-MVP",
        ],
        ["SurfConext login", "NRO", "Geen multi-user authenticatie in MVP"],
    ],
    row_height=0.34,
)

add_note(
    slide18,
    'Bron: "Evaluatietool Selectie Hoger Onderwijs" (NRO 2023), "Ontwerpfase ETHO" (Dialogic 2023), "Technisch ontwerp v1.0" (Yard 2023)',
    top=6.5,
)
fixes.append(
    "Slide 19: bronnen per rij, variabele-definitie uitgelegd, config-tool verwijderd (eigen slide), CSV weg"
)


# ============================================================
# Fix "oorspronkelijke opdracht" -> "oorspronkelijk idee" in twijfelgevallen slide
# ============================================================
for idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if "twijfelgevallen" in para.text.lower() or "Open vragen" in para.text:
                    # Found it - fix all text on this slide
                    for shape2 in slide.shapes:
                        if shape2.has_text_frame:
                            for para2 in shape2.text_frame.paragraphs:
                                for run in para2.runs:
                                    if "oorspronkelijke opdracht" in run.text:
                                        run.text = run.text.replace(
                                            "oorspronkelijke opdracht",
                                            "oorspronkelijke idee",
                                        )
                                        fixes.append(
                                            "Twijfelgevallen: opdracht -> idee"
                                        )
                    break


# ============================================================
# Fix CSV mention in slide 16 (over alle bestanden heen) - remove "alleen Excel"
# ============================================================
for shape in prs.slides[16].shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "oorspronkelijke opdracht" in run.text:
                    run.text = run.text.replace(
                        "oorspronkelijke opdracht", "oorspronkelijke idee"
                    )


try:
    prs.save("20260604_selectietool_mvp_voorstel_intern.pptx")
except PermissionError:
    prs.save("20260604_selectietool_mvp_voorstel_intern_v2.pptx")
    print("Opgeslagen als _v2 (origineel open in PowerPoint).")
print("Presentatie gefixt. %d fixes:" % len(fixes))
for f in fixes:
    print("  - " + f)
