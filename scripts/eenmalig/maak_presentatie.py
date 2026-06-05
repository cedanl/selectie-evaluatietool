"""
Genereer de presentatie voor de Evaluatietool Selectie.

Vertelt het verhaal van het project: het probleem, de aanpak, de technische
werking, en wat het dashboard laat zien. Inclusief config wizard en
PDF-rapportgeneratie.

Draai:
    uv run python scripts/maak_presentatie.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
DARK = RGBColor(51, 51, 51)
GRAY = RGBColor(120, 120, 120)
BLUE = RGBColor(44, 62, 80)
LIGHT_BG = RGBColor(245, 245, 245)
ACCENT = RGBColor(41, 128, 185)
GREEN = RGBColor(39, 174, 96)
ORANGE = RGBColor(230, 126, 34)


def add_slide(title_text, body_lines=None, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    title = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.9)
    )
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BLUE

    if body_lines:
        body = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2)
        )
        tf = body.text_frame
        tf.word_wrap = True
        for i, line in enumerate(body_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

            if line == "":
                p.space_before = Pt(8)
                continue

            bold = False
            if line.startswith("**") and line.endswith("**"):
                line = line[2:-2]
                bold = True

            if line.startswith("- "):
                p.text = line
                p.level = 0
                p.space_before = Pt(4)
            else:
                p.text = line
                if bold:
                    p.font.bold = True
                p.space_before = Pt(4)

            p.font.size = Pt(18)
            p.font.color.rgb = DARK

    if note:
        nb = slide.shapes.add_textbox(
            Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.6)
        )
        tf = nb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = note
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY

    return slide


def add_table_slide(title_text, headers, rows, intro=None, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    title = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.9)
    )
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BLUE

    top = Inches(1.5)
    if intro:
        ib = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.6)
        )
        tf = ib.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = intro
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
        top = Inches(2.2)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_width = Inches(11.7)
    tbl_height = Inches(0.4) * n_rows
    shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.8), top, tbl_width, tbl_height
    )
    table = shape.table

    col_width = int(tbl_width / n_cols)
    for ci in range(n_cols):
        table.columns[ci].width = col_width

    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = DARK
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG

    if note:
        nb = slide.shapes.add_textbox(
            Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.6)
        )
        tf = nb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = note
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY

    return slide


# ============================================================
# SLIDE 1: Titel
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = BLUE

title = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.5))
tf = title.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Evaluatietool Selectie"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.LEFT

sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(11.7), Inches(1.0))
tf = sub.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = (
    "Selectiedata koppelen aan studiesucces, ongeacht de opleiding. "
    "Een werkend dashboard met config wizard en PDF-rapport."
)
p.font.size = Pt(22)
p.font.color.rgb = RGBColor(200, 200, 200)

foot = slide.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.5))
tf = foot.text_frame
p = tf.paragraphs[0]
p.text = "CEDA / Tanjung Analytics"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(160, 160, 160)


# ============================================================
# SLIDE 2: De vraag
# ============================================================
add_slide(
    "De vraag die elke opleiding heeft",
    [
        "Universiteiten en hogescholen selecteren studenten voor opleidingen met beperkte plekken.",
        "Elke opleiding doet dat op een andere manier: de ene gebruikt toetsen, de andere kijkt naar "
        "cijfers, en weer een andere doet gesprekken.",
        "",
        "Na de selectie wil je weten: werkten onze selectiecriteria? Presteren de geselecteerde studenten "
        "beter dan de niet-geselecteerden?",
        "",
        "Het probleem: elk selectiebestand ziet er anders uit. Andere kolommen, andere schalen, "
        "andere structuur. Daardoor moet je voor elke opleiding opnieuw beginnen met je analyse.",
        "",
        "Deze tool lost dat op. De opleiding levert een selectiebestand en een configuratiebestand aan. "
        "De tool doet de rest: omzetten, koppelen, en een dashboard tonen.",
    ],
)


# ============================================================
# SLIDE 3: Bewijs dat het probleem echt is
# ============================================================
add_table_slide(
    "Elk selectiebestand is anders",
    [
        "Kenmerk",
        "FAR Leiden 2025",
        "FAR Leiden 2026",
        "Psychologie 2022-2023",
        "Psychologie 2026-2027",
    ],
    [
        ["Opleidingsniveau", "Master", "Master", "Bachelor", "Bachelor"],
        ["Aantal kolommen", "60", "97", "9", "46"],
        [
            "Selectiemethode",
            "Gesprekken + diploma",
            "Digitale toetsen",
            "Toets + matching + cijfers",
            "Schoolcijfers + matching",
        ],
        [
            "Menselijke beoordelaar?",
            "Ja, twee per kandidaat",
            "Alleen bij open vraag",
            "Nee",
            "Nee",
        ],
        [
            "ID-kolom",
            "A_Nummer_Aanvraag",
            "Studentnummer",
            "Studentnummer",
            "Studentnummer",
        ],
        [
            "Kolomnamen boven de data",
            "Rij 1",
            "Rij 1",
            "Rij 1",
            "Rij 3 (groepskoppen boven)",
        ],
    ],
    intro="Vier van de vijf bestanden waarmee de tool getest is. Elk ziet er compleet anders uit.",
)


# ============================================================
# SLIDE 4: De oplossing
# ============================================================
add_slide(
    "Hoe werkt de tool?",
    [
        "De gebruiker uploadt drie bestanden in het dashboard:",
        "",
        "**1. Selectiebestand (Excel)**",
        "Het bestand met de resultaten van de selectieprocedure. Dit verschilt per opleiding.",
        "",
        "**2. Configuratiebestand (Excel of via de config wizard)**",
        "Beschrijft welke kolommen uit het selectiebestand worden meegenomen en hoe ze heten. "
        "Kan automatisch gegenereerd worden met de ingebouwde config wizard.",
        "",
        "**3. 1CHO-data (studiesuccesdata)**",
        "Bevat de groepindeling per student: niet gestart, gestart maar niet doorgestroomd, "
        "of doorgestroomd naar jaar 2.",
        "",
        "Het dashboard opent zodra alle drie bestanden geladen zijn. "
        "Na de analyse kun je een PDF-rapport downloaden.",
    ],
)


# ============================================================
# SLIDE 5: Wat doet de analist?
# ============================================================
add_slide(
    "Wat doet de analist?",
    [
        "De analist is verantwoordelijk voor het selectiebestand en de configuratie.",
        "",
        "**Selectiebestand klaarzetten**",
        "De analist krijgt het selectiebestand van de opleiding of het testbureau. "
        "Het bestand hoeft niet opgeschoond te worden. Laat het zoals het is.",
        "",
        "**Configuratie maken**",
        "Open de tool en upload het selectiebestand. Klik op 'config automatisch genereren'. "
        "De wizard detecteert welke kolommen scores bevatten. De analist controleert het resultaat: "
        "kloppen de instrumentnamen? Zijn er kolommen die niet thuishoren? "
        "Pas aan waar nodig en klik 'Bevestig config'.",
        "",
        "**Inhoudelijke keuzes**",
        "De analist beslist welke kolommen meetellen. Bij meerdere scoreversies "
        "(bijvoorbeeld schaalscore en normscore): kies er een. Bij meerdere beoordelaars: "
        "gebruik de samengevoegde score. De wizard neemt alles mee, de analist filtert.",
    ],
)


# ============================================================
# SLIDE 6: Wat doet de dataprofessional?
# ============================================================
add_slide(
    "Wat doet de dataprofessional?",
    [
        "De dataprofessional is verantwoordelijk voor de 1CHO-data (studiesucces).",
        "",
        "**1CHO-data ophalen**",
        "De dataprofessional haalt de studiesuccesgegevens op uit de 1 Cijfer HO bestanden van DUO. "
        "Hiervoor kan het 1cijferho-project van CEDA gebruikt worden (github.com/cedanl/1cijferho). "
        "Dat project heeft een preset voor de selectietool die precies de juiste kolommen selecteert.",
        "",
        "**Filteren op opleiding**",
        "Filter de 1CHO-data op de opleiding die je analyseert. Het resultaat bevat per student: "
        "studentnummer, selectiejaar, en de groepsindeling (niet gestart, uitval, of doorgestroomd). "
        "Optioneel ook geslacht, herkomst, vooropleiding en VO-eindcijfer.",
        "",
        "**Studentnummer checken**",
        "Controleer dat het studentnummer in hetzelfde formaat staat als in het selectiebestand. "
        "Let op voorloopnullen: als het ene bestand 0012345 gebruikt en het andere 12345, "
        "dan mislukt de koppeling.",
    ],
)


# ============================================================
# SLIDE 7: Het configuratiebestand
# ============================================================
add_slide(
    "Het configuratiebestand: het scharnierpunt",
    [
        "Per selectiebestand heb je een configuratiebestand nodig. Er zijn twee manieren om er een te maken:",
        "",
        "**Config wizard (aanbevolen)**",
        "De ingebouwde wizard leest het selectiebestand en detecteert automatisch welk blad de data bevat, "
        "waar de kolomnamen staan, welke kolom het studentnummer is, en welke kolommen scores bevatten. "
        "De gebruiker controleert het resultaat in een bewerkbare tabel en past namen aan waar nodig.",
        "",
        "**Handmatig via het templatebestand**",
        "Een Excel-spreadsheet met twee tabbladen: Instellingen (studentnummer, headerrij, totaalscore) "
        "en Kolommen (per scorekolom een instrument, item, criterium en score-type).",
        "",
        "De tool schrijft niks voor. De analist kiest zelf de namen voor instrumenten, items en criteria. "
        "Dat maakt het flexibel genoeg om voor elke opleiding te werken.",
    ],
)


# ============================================================
# SLIDE 6: Voorbeeld config kolommen
# ============================================================
add_table_slide(
    "Voorbeeld: kolommen-tabblad (Farmacie 2026)",
    ["kolom_naam", "instrument", "item", "criterium", "score_type"],
    [
        [
            "ctb_reflecteren_Schaalscore",
            "Competentietest",
            "Reflecteren schaalscore",
            "Reflectievermogen",
            "schaalscore",
        ],
        [
            "ctb_stressbestendigheid_Schaalscore",
            "Competentietest",
            "Stressbestendigheid schaalscore",
            "Stressbestendigheid",
            "schaalscore",
        ],
        [
            "Beoordeling_reflectievermogen...",
            "Open vraag",
            "Beoordeling reflectie (1-2-3)",
            "Reflectievermogen",
            "schaal 1-3",
        ],
        [
            "sjts_totaal_Schaalscore",
            "SIT-S",
            "Totaalscore sociale intelligentie",
            "Sociale intelligentie",
            "schaalscore",
        ],
    ],
    intro="Dit bestand heeft 97 kolommen. De config selecteert er 12 en geeft ze leesbare namen.",
    note="De analist kiest zelf de namen voor instrument, item en criterium. Er is geen vaste lijst.",
)


# ============================================================
# SLIDE 7: Voorbeeld config kolommen Psychologie
# ============================================================
add_table_slide(
    "Ander voorbeeld: kolommen-tabblad (Psychologie 2026-2027)",
    ["kolom_naam", "instrument", "item", "criterium", "score_type"],
    [
        [
            "WI SCORE",
            "Schooldiploma",
            "Wiskunde puntenscore",
            "Vakkennis wiskunde",
            "punten (0-5)",
        ],
        [
            "EN SCORE",
            "Schooldiploma",
            "Engels puntenscore",
            "Vakkennis Engels",
            "punten (0-5)",
        ],
        [
            "V1 SCORE",
            "Schooldiploma keuzevak",
            "Keuzevak 1 puntenscore",
            "Vakkennis keuzevak",
            "punten (0-5)",
        ],
        [
            "Matching Score (1-3)",
            "Matchingsvragenlijst",
            "Matchingscore (1-3)",
            "Studiemotivatie",
            "schaal 1-3",
        ],
        [
            "Vragenlijst %",
            "Deelscore",
            "Vragenlijst score percentage",
            "",
            "percentage",
        ],
    ],
    intro="Compleet andere structuur: schoolcijfers, keuzevakken en matchingscore. Dezelfde config-structuur werkt.",
    note="Van de 46 kolommen zijn er 19 in de config opgenomen.",
)


# ============================================================
# SLIDE 8: Het standaardformaat
# ============================================================
add_table_slide(
    "Het standaardformaat: elke score op een eigen rij",
    [
        "studentnummer",
        "opleiding",
        "jaar",
        "instrument",
        "item",
        "criterium",
        "score",
    ],
    [
        [
            "123456",
            "Farmacie",
            "2026",
            "Competentietest",
            "Reflecteren schaalscore",
            "Reflectievermogen",
            "30",
        ],
        [
            "123456",
            "Farmacie",
            "2026",
            "SIT-S",
            "Totaalscore sociale intelligentie",
            "Sociale intelligentie",
            "51",
        ],
        [
            "9876543",
            "Psychologie",
            "2026",
            "Schooldiploma",
            "Wiskunde puntenscore",
            "Vakkennis wiskunde",
            "3.80",
        ],
        [
            "9876543",
            "Psychologie",
            "2026",
            "Matchingsvragenlijst",
            "Matchingscore (1-3)",
            "Studiemotivatie",
            "3",
        ],
    ],
    intro="De tool zet elk selectiebestand om naar dit formaat. Elke rij is een meting van een kandidaat op een item.",
    note="Totaalscore en rangnummer staan in de instellingen-tab van de config, niet in het kolommen-tabblad.",
)


# ============================================================
# SLIDE 9: Drie groepen
# ============================================================
add_slide(
    "Drie groepen: de kern van de analyse",
    [
        "De tool vergelijkt selectiescores tussen drie groepen kandidaten. De groepindeling komt uit "
        "de 1CHO-studiesuccesdata.",
        "",
        "**Niet gestart**",
        "Kandidaten die niet in 1CHO staan. Ze zijn niet toegelaten, of zijn wel geselecteerd maar "
        "nooit gestart met de opleiding.",
        "",
        "**Gestart, niet naar jaar 2**",
        "Studenten die wel zijn gestart maar niet zijn doorgestroomd naar het tweede jaar. "
        "Uitval of overstap na het eerste jaar.",
        "",
        "**Doorgestroomd naar jaar 2**",
        "Studenten die het eerste jaar hebben volbracht en doorgaan.",
        "",
        "Per selectie-item kun je dan zien: scoren de drie groepen anders? "
        "Als de doorstromers hoger scoren dan de uitvallers, heeft het selectie-instrument "
        "voorspellende waarde.",
    ],
)


# ============================================================
# SLIDE 10: Dashboard tab 1 - Selectiescores
# ============================================================
add_slide(
    "Dashboard: Selectiescores per groep",
    [
        "Het eerste tabblad toont per selectie-item een boxplot met de drie groepen naast elkaar.",
        "",
        "**Wat je hier ziet**",
        "- Mediaan, spreiding en individuele scores per groep",
        "- Filters op instrument, criterium, cohort, geslacht en vooropleiding",
        "- Tabel met gemiddelden en standaarddeviaties per groep",
        "",
        "**Wat je eruit kunt halen**",
        "Als de groene boxen (doorgestroomd) consequent hoger liggen dan de oranje (uitval), "
        "heeft het item voorspellende waarde.",
        "Als de boxen grotendeels overlappen, meet het item niet wat je hoopt te meten.",
        "",
        "**Belangrijk**",
        "De scores zijn niet genormaliseerd. Je vergelijkt alleen binnen een item, niet tussen items. "
        "Een schaalscore van 25 en een punten-score van 3 zijn niet vergelijkbaar.",
    ],
)


# ============================================================
# SLIDE 11: Dashboard tab 2 - Samenhang
# ============================================================
add_slide(
    "Dashboard: Samenhang en regressie",
    [
        "Het tweede tabblad combineert twee analyses: hoe hangen items onderling samen, en welke "
        "items voorspellen studiesucces?",
        "",
        "**Correlatiematrix**",
        "Een heatmap die toont hoe sterk elk paar items samenhangt. Hoge correlaties betekenen dat "
        "twee items grotendeels hetzelfde meten. Items die weinig correleren voegen elk unieke "
        "informatie toe aan de selectie.",
        "",
        "**Logistische regressie**",
        "Een model dat voorspelt wie doorstroomt naar jaar 2, op basis van alle selectie-items. "
        "Per item krijg je een coefficient, odds ratio en p-waarde.",
        "",
        "Een significant item (p < 0.05) voegt voorspellende waarde toe, ook als je de andere items "
        "al meeneemt. Een niet-significant item voegt weinig toe boven wat de andere items al verklaren.",
        "",
        "Het model toont ook de pseudo R-kwadraat: hoeveel verklaart het geheel aan items samen?",
    ],
)


# ============================================================
# SLIDE 12: Uitdagingen bij de regressie
# ============================================================
add_slide(
    "Uitdagingen bij de logistische regressie",
    [
        "Selectiedata is lastig voor regressie. Vier problemen komen steeds terug:",
        "",
        "**Verschillende schalen**",
        "Het ene item is een schaalscore van 0-100, het andere een beoordeling van 1-3. "
        "In een rauwe regressie domineren items op grotere schalen het model, puur door hun bereik.",
        "",
        "**Weinig studenten, veel items**",
        "Een selectie heeft vaak 50-150 ingeschreven studenten en 10+ items. "
        "De vuistregel is 5 studenten in de kleinste groep per predictor. Met 15 uitvallers "
        "kun je maar 3 items tegelijk meenemen. Meer predictoren geeft instabiele resultaten.",
        "",
        "**Overlap tussen items**",
        "Selectie-instrumenten meten vaak vergelijkbare dingen. Twee subschalen van dezelfde "
        "competentietest correleren al snel r = 0.7. In een gezamenlijk model wordt dan geen "
        "van beide significant, terwijl ze individueel wel voorspellend zijn.",
        "",
        "**Ontbrekende data**",
        "Niet elke kandidaat maakt elk onderdeel. Keuzevakken, optionele toetsen en afgebroken "
        "testen zorgen voor gaten in de data.",
    ],
)


# ============================================================
# SLIDE 13: Hoe de tool hiermee omgaat
# ============================================================
add_slide(
    "Hoe de tool hiermee omgaat",
    [
        "**Schaalverschillen: z-score normalisatie**",
        "Alle scores worden omgerekend naar z-scores (gemiddelde = 0, SD = 1) voor de regressie. "
        "Daardoor drukken de coefficienten en odds ratios het effect uit per standaarddeviatie, "
        "niet per ruwe punt. Dat maakt items op verschillende schalen vergelijkbaar.",
        "",
        "**Te weinig studenten: automatische selectie**",
        "De tool berekent hoeveel predictoren het model aankan (kleinste groep / 5). "
        "Als er meer items zijn, worden ze eerst individueel getoetst. Alleen de sterkste "
        "items gaan het gezamenlijke model in. De rest wordt gerapporteerd als 'niet meegenomen'.",
        "",
        "**Overlap: multicollineariteitscheck**",
        "Voor het fitten controleert de tool of de predictormatrix vol rang is. "
        "Items die lineair afhankelijk zijn van andere items worden verwijderd. "
        "De correlatiematrix op hetzelfde tabblad helpt de gebruiker overlap te herkennen.",
        "",
        "**Ontbrekende data: drempel + imputatie**",
        "Items met meer dan 30% ontbrekende waarden worden uitgesloten. "
        "Bij de rest worden missende waarden opgevuld met het kolomgemiddelde.",
    ],
)


# ============================================================
# SLIDE 14: Dashboard tab 3 - Demografisch
# ============================================================
add_slide(
    "Dashboard: Demografisch profiel",
    [
        "Het derde tabblad toont de verdeling van achtergrondkenmerken per groep.",
        "",
        "**Verdeling per cohort**",
        "Een gestapeld staafdiagram per selectiejaar: hoeveel procent valt in elke groep? "
        "Hiermee zie je of het selectieproces van jaar tot jaar verandert.",
        "",
        "**Geslacht, herkomst en vooropleiding**",
        "Per groep: hoe ziet de samenstelling eruit? Zijn er systematische verschillen in wie "
        "doorstroomt en wie uitvalt?",
        "",
        "De achtergrondkenmerken komen uit de 1CHO-data en zijn alleen beschikbaar voor studenten "
        "die daadwerkelijk ingeschreven zijn geweest.",
        "",
        "Dit is geen fairness-analyse. Het toont patronen, maar trekt geen conclusies over "
        "of de selectie eerlijk verloopt. Dat vraagt om aanvullende methodiek.",
    ],
)


# ============================================================
# SLIDE 13: Dashboard tab 4 - VO-cijfer
# ============================================================
add_slide(
    "Dashboard: VO-eindcijfer vs selectiescores",
    [
        "Het vierde tabblad vergelijkt selectiescores met het gemiddelde VO-eindcijfer van elke student.",
        "",
        "**Het idee**",
        "Het VO-eindcijfer is een onafhankelijke meting van cognitieve prestaties. Het is niet "
        "onderdeel van de selectie, maar komt uit 1CHO.",
        "",
        "**Lage correlatie (r rond 0)**",
        "Het selectie-item meet iets anders dan schoolprestaties. Dat is vaak gewenst: de selectie "
        "voegt dan informatie toe die het VO-cijfer niet al geeft.",
        "",
        "**Hoge correlatie (r > 0.5)**",
        "Het selectie-item overlapt sterk met het VO-cijfer. De selectie herhaalt dan wat het "
        "schooldiploma al vertelt.",
        "",
        "De tabel toont Pearson r voor elk item en de totaalscore. Per item is er een scatterplot "
        "met trendlijnen per groep.",
    ],
)


# ============================================================
# SLIDE 14: Wat we bewust overslaan
# ============================================================
add_table_slide(
    "Wat we per bestand meenemen en overslaan",
    [
        "Bestand",
        "Kolommen totaal",
        "Kolommen in config",
        "Overgeslagen",
        "Voornaamste reden",
    ],
    [
        [
            "FAR Leiden 2025",
            "60",
            "11",
            "49",
            "Dubbele beoordelaars, tekstvelden, persoonsgegevens",
        ],
        [
            "FAR Leiden 2026",
            "97",
            "12",
            "85",
            "Drie scoreversies, proctoring, procesdata",
        ],
        [
            "Psychologie 2022-2023",
            "9",
            "6",
            "3",
            "Studentnummer, totaalscore, rangnummer (uitkomsten)",
        ],
        [
            "Psychologie 2026-2027",
            "46",
            "19",
            "27",
            "Vaknamen, ruwe cijfers, herhaalde structuur",
        ],
    ],
    intro="De config bepaalt wat meekomt. De rest wordt genegeerd. Liever bewust selecteren dan alles meenemen.",
)


# ============================================================
# SLIDE 15: Keuzes per bestand
# ============================================================
add_slide(
    "Keuzes die de analist maakt in de config",
    [
        "Elk selectiebestand vraagt om inhoudelijke keuzes. De config wizard helpt met detectie, "
        "maar de analist controleert altijd het resultaat.",
        "",
        "**Bij meerdere scoreversies: kies er een**",
        "Farmacie 2026 heeft per toets een schaalscore, normscore en Z-score. "
        "De wizard neemt alles mee. De analist verwijdert de overbodige rijen uit de tabel.",
        "",
        "**Bij meerdere beoordelaars: gebruik de samengevoegde score**",
        "Farmacie 2025 heeft twee beoordelaars per kandidaat. "
        "De analist verwijdert de individuele B1/B2 rijen en houdt de samengevoegde C_-scores over.",
        "",
        "**Bij extra koprijen: stel de header_rij in**",
        "Psychologie 2026-2027 heeft groepskoppen boven de kolomnamen. "
        "De wizard probeert de juiste rij te detecteren, maar de analist kan het handmatig aanpassen.",
        "",
        "**Tekstvelden, datums en persoonsgegevens**",
        "De wizard filtert deze automatisch uit. Alleen numerieke kolommen worden voorgesteld.",
    ],
)


# ============================================================
# SLIDE 16: Drie begrippen
# ============================================================
add_table_slide(
    "Drie begrippen die de config gebruikt",
    ["Begrip", "Wat is het?", "Voorbeeld"],
    [
        [
            "Instrument",
            "Het middel waarmee je beoordeelt: een toets, een gesprek, een vragenlijst.",
            "Competentietest, Schooldiploma, Gesprek",
        ],
        [
            "Item",
            "Een onderdeel van een instrument. Dat kan een losse meting zijn, "
            "of een apart beoordelingspunt.",
            "Reflecteren schaalscore, Wiskunde puntenscore",
        ],
        [
            "Criterium",
            "De eigenschap die je wilt meten. Een item meet altijd (optioneel) een criterium.",
            "Reflectievermogen, Vakkennis wiskunde",
        ],
    ],
    intro="De config beschrijft elke kolom op drie niveaus: van breed (instrument) naar smal (criterium).",
    note="De analist kiest zelf welke namen bij een kolom horen. De tool schrijft niks voor.",
)


# ============================================================
# SLIDE 17: Wat we niet doen
# ============================================================
add_slide(
    "Wat de tool niet doet (en waarom)",
    [
        "**Geen normalisatie van scores**",
        "Schaalscores, percentages en 1-2-3-beoordelingen zijn niet vergelijkbaar. "
        "We vergelijken alleen binnen een item, niet tussen items.",
        "",
        "**Geen fairness-analyse**",
        "De demografische tab toont patronen, maar de tool trekt geen conclusies over eerlijkheid. "
        "Dat vraagt om aanvullende methodiek en zorgvuldige interpretatie.",
        "",
        "**Geen multi-user authenticatie**",
        "De tool draait nu lokaal. SurfConext login en rolbeheer zijn post-scope.",
        "",
        "**Geen multi-bestandsupload per selectie**",
        "Sommige opleidingen hebben meerdere bestanden per selectieprocedure. "
        "De tool ondersteunt nu een selectiebestand per keer.",
        "",
        "**Geen cohortanalyse over meerdere jaren**",
        "Je analyseert per selectiejaar. Vergelijking tussen jaargangen is een volgende stap.",
    ],
)


# ============================================================
# SLIDE 19: Relatie tot oorspronkelijke opdracht
# ============================================================
add_table_slide(
    "Dekking van de oorspronkelijke opdracht",
    ["Onderdeel", "Status", "Toelichting"],
    [
        [
            "Selectiedata uploaden en koppelen",
            "Gedaan",
            "Config-gebaseerd, getest met 4 bestanden",
        ],
        [
            "Dashboard met groepsvergelijking",
            "Gedaan",
            "Boxplots per item, drie groepen",
        ],
        [
            "Koppeling met 1CHO-studiesuccesdata",
            "Gedaan",
            "Via studentnummer, automatische groepindeling",
        ],
        [
            "Beschrijvende statistiek per item",
            "Gedaan",
            "Gemiddelden, SD, spreiding per groep",
        ],
        [
            "Correlatie- en regressieanalyse",
            "Gedaan",
            "Correlatiematrix + logistische regressie",
        ],
        [
            "Demografisch profiel per groep",
            "Gedaan",
            "Geslacht, herkomst, vooropleiding uit 1CHO",
        ],
        [
            "VO-eindcijfer vs selectiescores",
            "Gedaan",
            "Pearson r en scatterplots",
        ],
        [
            "Config wizard (auto-detectie kolommen)",
            "Gedaan",
            "Detecteert blad, header, ID, scores automatisch",
        ],
        [
            "PDF-evaluatierapport",
            "Gedaan",
            "Downloadbaar rapport met grafieken en tabellen",
        ],
        [
            "Fairness-analyse",
            "Niet gedaan",
            "Patronen zichtbaar, geen conclusies",
        ],
        [
            "Webapplicatie met authenticatie",
            "Niet gedaan",
            "Draait lokaal",
        ],
    ],
    intro="De oorspronkelijke opdrachtomschrijving (NRO, 2023) beschrijft een complete webapplicatie. "
    "Dit is wat we daarvan gebouwd hebben.",
)


# ============================================================
# SLIDE 20: Volgende stappen
# ============================================================
add_slide(
    "Volgende stappen",
    [
        "**Korte termijn**",
        "- Meer opleidingen testen met hun eigen data",
        "- Feedback verwerken op dashboard-layout en analyses",
        "- Tekstuele duiding toevoegen aan het PDF-rapport",
        "",
        "**Middellange termijn**",
        "- Cohortanalyse: meerdere jaargangen vergelijken binnen een opleiding",
        "- Effectgroottes (Cohen's d) naast p-waarden in de regressie",
        "- Interrater-analyse voor bestanden met meerdere beoordelaars",
        "- Meerdere selectiebestanden per opleiding ondersteunen",
        "",
        "**Langere termijn**",
        "- Fairness-analyse met gevalideerde methodiek",
        "- Deployment als webapplicatie met SurfConext authenticatie",
        "- Integratie met 1cijferho-tool voor directe data-import",
    ],
)


# ============================================================
# SLIDE 21: Samenvatting
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = BLUE

title = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.0))
tf = title.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Samenvatting"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE

body = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(11.7), Inches(3.0))
tf = body.text_frame
tf.word_wrap = True

lines = [
    "Elk selectiebestand is anders. De tool lost dat op met een configuratiebestand "
    "dat je handmatig maakt of automatisch laat genereren via de config wizard.",
    "",
    "Getest met vier bestanden van twee opleidingen: Farmacie (master, LUMC) en Psychologie (bachelor, Leiden).",
    "",
    "Het dashboard toont selectiescores per groep, correlaties, regressie, demografie en VO-eindcijfer. "
    "Na de analyse kun je een PDF-rapport downloaden.",
    "",
    "De volgende stap: meer opleidingen laten testen met hun eigen data.",
]

for i, line in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    if line == "":
        p.space_before = Pt(10)
        continue
    p.text = line
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.space_before = Pt(6)


prs.save("20260604_selectietool_mvp_voorstel_intern.pptx")
print("Presentatie opgeslagen als 20260604_selectietool_mvp_voorstel_intern.pptx")
